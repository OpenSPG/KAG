from pptx import Presentation
from pptx.util import Pt, Inches
import re


def create_slide(prs, title, content=None, content_level=None):
    """
    创建幻灯片，并根据标题和内容添加到幻灯片上。
    参数:
        prs: Presentation 对象
        title: 幻灯片标题
        content: 幻灯片正文内容
        content_level: 内容的层级 (用于调整缩进)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 用第一个母版生成一页ppt
    for shape in slide.placeholders:         # 获取这一页所有的占位符
        phf = shape.placeholder_format
        print(f'{phf.idx}--{shape.name}--{phf.type}')
        print('shape name ', shape.name)
        if shape.name ==  'Title 1':
            shape.text = f'目标'  #在标题占位符中填写“目标”
        else:
            shape.text = f'内容'  #在其他占位符中填写“内容”

    # slide = prs.slides.add_slide(prs.slide_layouts[1])  # 使用标题+正文布局
    # title_placeholder = slide.shapes.title
    # content_placeholder = slide.placeholders[1]  # 获取正文框
    #
    # title_placeholder.text = title  # 设置标题
    #
    # if content:
    #     # 按内容层级缩进
    #     p = content_placeholder.text_frame.add_paragraph()
    #     p.text = content
    #     if content_level == 1:
    #         p.level = 1  # 一级缩进
    #     elif content_level == 2:
    #         p.level = 2  # 二级缩进
    #     elif content_level == 3:
    #         p.level = 3  # 三级缩进
    #     else:
    #         p.level = 0  # 默认无缩进


def parse_text_to_ppt(text, output_file="output.pptx"):
    """
    解析文本并生成 PPT 文件。
    参数:
        text: 输入的多层级文本
        output_file: 输出的 PPT 文件名
    """
    prs = Presentation()  # 创建一个空的 PPT

    current_slide_title = None
    for line in text.split("\n"):
        line = line.strip()

        # 跳过空行
        if not line:
            continue

        # 检测标题层级
        if line.startswith("####"):  # 四级标题
            match = re.match(r"####\s*(.+)", line)
            if match:
                create_slide(prs, current_slide_title, content=match.group(1), content_level=3)
        elif line.startswith("###"):  # 三级标题
            match = re.match(r"###\s*(.+)", line)
            if match:
                create_slide(prs, current_slide_title, content=match.group(1), content_level=2)
        elif line.startswith("##"):  # 二级标题
            match = re.match(r"##\s*(.+)", line)
            if match:
                current_slide_title = match.group(1)
                create_slide(prs, current_slide_title)
        elif line.startswith("#"):  # 一级标题
            match = re.match(r"#\s*(.+)", line)
            if match:
                current_slide_title = match.group(1)
                create_slide(prs, current_slide_title)
        else:
            # 普通正文内容
            create_slide(prs, current_slide_title, content=line, content_level=1)

    # 保存最终的 PPT
    prs.save(output_file)


if __name__ == "__main__":
    with open('../doc.txt', 'r', encoding='utf-8') as file:
        content = file.read()  # 读取整个文件内容
    # 调用生成 PPT
    parse_text_to_ppt(content, output_file="立体纸模服装制作.pptx")