import re
from typing import List
from pptx.util import Inches
from typing import List, Tuple
from bs4 import BeautifulSoup
from pptx.util import Inches, Pt


class TablePPTCreator:
    def __init__(self, prs):
        """接收统一的 Presentation 对象"""
        self.prs = prs

    def add_table_slide(self, title, content, rows: int, cols: int, table_content: List[List[str]], position=None):
        """添加表格页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[8])

        if slide.shapes.title:
            # 如果存在标题占位符，设置标题文本
            title_shape = slide.shapes.title
            title_shape.text = title
        if content is not None:
            shape = slide.shapes[1]
            shape.text = content
        # 设置表格位置和大小
        shape = slide.shapes[2]
        placeholder = shape
        table_left, table_top, table_width, table_height = (
            placeholder.left,
            placeholder.top,
            placeholder.width,
            placeholder.height,
        )
        # 添加表格
        table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

        for i in range(rows):
            for j in range(cols):
                # 设置单元格内容
                cell = table.cell(i, j)
                cell.text = table_content[i][j]

        return slide

    def add_html_table_slide_with_merge(
            self, title: str, content, table_content: List[List[str]], merge_cells: List[Tuple[int, int, int, int]]
    ):
        """
        添加带有合并单元格的表格到 PowerPoint 幻灯片中
        :param content: 对表格的解释内容，无则填写None
        :param title: 幻灯片标题
        :param table_content: 表格内容，二维列表
        :param merge_cells: 合并单元格信息 (start_row, start_col, end_row, end_col)
        """
        # 创建一个空白布局幻灯片
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[8])
        # 设置标题
        if slide.shapes.title:
            # 如果存在标题占位符，设置标题文本
            title_shape = slide.shapes.title
            title_shape.text = title

        if content is not None:
            shape = slide.shapes[1]
            shape.text = content

        rows = len(table_content)
        cols = max(len(row) for row in table_content)
        # 设置表格位置和大小
        shape = slide.shapes[2]
        placeholder = shape
        table_left, table_top, table_width, table_height = (
            placeholder.left,
            placeholder.top,
            placeholder.width,
            placeholder.height,
        )

        # 添加表格
        table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

        # 填充内容并合并单元格
        for (start_row, start_col, end_row, end_col) in merge_cells:
            try:
                # 合并单元格
                table.cell(start_row, start_col).merge(
                    table.cell(end_row, end_col))
            except IndexError:
                # 添加错误处理
                print(f"警告：无效的合并范围 ({start_row},{start_col})-({end_row},{end_col})")

        # 填充文本内容
        for row_idx, row in enumerate(table_content):
            for col_idx, text in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text_frame.paragraphs[0].text = text
                cell.text_frame.paragraphs[0].font.size = Pt(12)

        return slide

    def add_markdown_table_slide(self, title, content, table_data):
        """
        在 PowerPoint 中添加一个表格幻灯片，支持样式设置
        :param content: 表格的解释内容，无则填写None
        :param title: 幻灯片标题
        :param table_data: 二维列表形式的表格内容
        """
        # 创建一个空白布局幻灯片
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[8])
        # 设置标题
        if slide.shapes.title:
            # 如果存在标题占位符，设置标题文本
            title_shape = slide.shapes.title
            title_shape.text = title

        if content is not None:
            shape = slide.shapes[1]
            shape.text = content

        rows = len(table_data)
        cols = max(len(row) for row in table_data)
        # 设置表格位置和大小
        shape = slide.shapes[2]
        placeholder = shape
        table_left, table_top, table_width, table_height = (
            placeholder.left,
            placeholder.top,
            placeholder.width,
            placeholder.height,
        )
        table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

        # 填充表格内容并设置样式
        for i, row in enumerate(table_data):
            for j, cell_text in enumerate(row):
                # 检查是否需要加粗（以 Markdown 加粗语法为依据）
                if "**" in cell_text:
                    text = re.sub(r"\*\*(.*?)\*\*", r"\1", cell_text)  # 去掉 Markdown 加粗语法
                    cell = table.cell(i, j)
                    cell.text = text
                    # 清除默认段落
                    cell.text_frame.clear()

                    # 添加新的段落并设置加粗
                    paragraph = cell.text_frame.add_paragraph()
                    run = paragraph.add_run()
                    run.text = text
                    run.font.bold = True
                else:
                    table.cell(i, j).text = cell_text
                table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(12)

    def save(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)
