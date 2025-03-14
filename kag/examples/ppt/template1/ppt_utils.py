import json
import os
import shutil
from copy import deepcopy

from PIL import Image
from pptx import Presentation
from typing import List, Dict
from bs4 import BeautifulSoup
import re
import re
from pathlib import Path


class PPTUtils:
    @staticmethod
    def check_font(font_name):
        """检查字体是否可用"""
        from matplotlib.font_manager import FontManager
        fm = FontManager()
        font_list = [f.name for f in fm.ttflist]
        return font_name in font_list

    @staticmethod
    def merge_ppts(ppt_list, output_path):
        """合并多个PPT"""
        merged_ppt = Presentation()

        for ppt_path in ppt_list:
            ppt = Presentation(ppt_path)

            for slide in ppt.slides:
                # 复制页面
                new_slide = merged_ppt.slides.add_slide(
                    merged_ppt.slide_layouts[
                        slide.slide_layout.index
                    ]
                )

                # 复制所有形状
                for shape in slide.shapes:
                    el = shape.element
                    new_el = deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(
                        new_el,
                        'p:extLst'
                    )

        merged_ppt.save(output_path)

    @staticmethod
    def copy_image_to_temp(image_path, target_dir=''):
        """复制图片到临时文件夹，并确保临时文件夹存在"""
        # 确定临时目录和目标路径
        os.makedirs(target_dir, exist_ok=True)  # 确保临时文件夹存在
        target_path = os.path.join(target_dir, os.path.basename(image_path))

        # 检查原始图片文件是否存在
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"图片文件未找到: {image_path}")

        # 使用 Pillow 打开图片，验证图片是否可用（可选）
        img = Image.open(image_path)
        img.verify()  # 验证图片完整性

        # 将图片复制到临时目录（直接复制文件）
        shutil.copy(image_path, target_path)

        print(f"图片已复制到: {target_path}")
        return target_path

    @staticmethod
    def get_placeholder_information(slide):
        """查看ppt母版的占位符信息，方便根据母版的对应占位符添加相关信息"""
        for placeholder in slide.placeholders:
            print(f"Placeholder {placeholder.placeholder_format.idx}: {placeholder.name}")

    @staticmethod
    def parse_text_to_json(text_file, output_file="output.json"):
        with open(text_file, "r", encoding="utf-8") as file:
            content = file.read()

        lines = content.strip().split("\n")  # 将文本按行分割
        json_data = {
            "title": None,
            "attributes": [],
            "pages": []
        }
        current_page = {}  # 当前处理的页面信息
        for line in lines:
            line = line.strip()  # 去除行首尾的空格
            if line.startswith("# "):  # 顶级标题（文件主标题）
                json_data["title"] = line[2:].strip()
            elif line.startswith("## "):  # 二级标题（目录属性）
                json_data["attributes"].append(line[3:].strip())
            elif line.startswith("### "):  # 三级标题（内容页标题）
                if current_page:  # 如果当前页面不为空，将当前页面信息存入 JSON
                    json_data["pages"].append(current_page)
                current_page = {"title": line[4:].strip(), "content": []}
            else:  # 内容部分（正文）
                if current_page:
                    current_page["content"].append(line)
        # 把最后一个页面添加进 JSON
        if current_page:
            json_data["pages"].append(current_page)

        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(json_data, f, ensure_ascii=False, indent=4)

        return json_data

    @staticmethod
    def extract_markdown_images(content):
        """
        提取Markdown文件中的图片信息
        :return: 包含图片信息的字典列表
        """
        # content = Path(md_file).read_text(encoding='utf-8')
        lines = content.split('\n')

        images = []

        for i, line in enumerate(lines):
            img_info = {}

            # 匹配HTML格式图片
            if '<img' in line:
                # 提取src属性
                src_match = re.search(r'src=["\'](.*?)["\']', line)
                if src_match:
                    img_info['url'] = src_match.group(1)

                    # 向前搜索标题和figureText
                    j = i - 1
                    while j >= 0:
                        prev_line = lines[j].strip()

                        # 匹配标题（非空行）
                        if prev_line and not prev_line.startswith(('<', '![', '|')):
                            img_info['title'] = prev_line
                            break
                        j -= 1

                    images.append(img_info)

            # 匹配Markdown标准图片
            elif line.startswith('!['):
                md_match = re.match(r'!\[(.*?)\]\((.*?)\)', line)
                if md_match:
                    img_info = {
                        'title': md_match.group(1),
                        'url': md_match.group(2)
                    }
                    images.append(img_info)

        return images

    @staticmethod
    def extract_markdown_tables(markdown_text):
        """
        从 Markdown 文件中提取所有表格，并解析为多个二维数组
        :return: List of tables, 每个表格为二维数组
        """
        # with open(md_file_path, "r", encoding="utf-8") as file:
        #     markdown_text = file.read()

        # 定义 Markdown 表格块的正则表达式
        table_block_regex = re.compile(
            r"(?:\|.*?\|(?:\n|\r|\r\n))+"
        )  # 匹配表格块（按行以 `|` 开头，并重复至少一行）

        # 查找所有表格块
        table_blocks = table_block_regex.findall(markdown_text)

        # 如果没有匹配到任何表格，返回空列表
        if not table_blocks:
            return []

        # 存储解析后的所有表格
        all_parsed_tables = []

        # 逐个表格块进行解析
        for table_block in table_blocks:
            # 按行分割表格块
            rows = table_block.strip().splitlines()

            # 移除分隔行（例如 "----"）
            parsed_table = []
            header_processed = False  # 标志是否已经处理表头
            for row in rows:
                # 按列分隔，并去掉多余空格
                columns = [col.strip() for col in row.strip('|').split('|')]

                # 判断是否为表头分隔符（例如 "----"）
                if not header_processed:
                    if all(re.match(r"^-+$", col) for col in columns):
                        header_processed = True  # 标记表头分隔符已处理
                        continue  # 跳过分隔符行
                    else:
                        # 标记表头内容加粗
                        columns = [f"**{col}**" for col in columns]

                parsed_table.append(columns)

            # 补齐列数不一致的情况
            max_cols = max(len(row) for row in parsed_table)
            for row in parsed_table:
                while len(row) < max_cols:
                    row.append("")  # 填充空单元格

            # 添加解析后的表格到结果列表
            all_parsed_tables.append(parsed_table)

        return all_parsed_tables

    @staticmethod
    def parse_markdown(markdown_file):
        """
        解析 Markdown 文件，包括一级标题（#）作为标题页，按二级标题（##）分块，
        并将所有三级标题（###）的内容聚合到对应二级标题的 content 中。
        :param markdown_file: Markdown 文件路径
        :return: [{'type': 'title', 'title': '一级标题内容'}, ...]
        """
        with open(markdown_file, 'r', encoding='utf-8') as f:
            markdown_text = f.read()

        # 用于存储解析后的结果
        sections = []

        # 提取一级标题（标题页）
        title_match = re.match(r'^#\s*(.+)', markdown_text, re.MULTILINE)
        if title_match:
            # 提取一级标题并存入 sections
            sections.append({
                'type': 'title',
                'title': title_match.group(1).strip(),
                'content': ''  # 标题页无内容
            })
            # 去掉一级标题部分，便于后续解析二级标题
            markdown_text = markdown_text[markdown_text.find('\n') + 1:]

        # 初始化变量
        current_section = None

        # 按行解析 Markdown 文本
        lines = markdown_text.splitlines()
        for line in lines:
            if line.startswith('## '):  # 二级标题
                # 如果有当前 section，将其添加到 sections 中
                if current_section:
                    sections.append(current_section)
                # 创建新的 section
                current_section = {
                    'type': 'section',
                    'title': line[3:].strip(),  # 去掉 ## 和空格
                    'content': ''
                }
            elif line.startswith('### '):  # 三级标题
                # 将三级标题直接追加到当前 section 的 content
                if current_section:
                    current_section['content'] += f"\n{line.strip()}\n"
            else:
                # 将普通内容直接追加到当前 section 的 content
                if current_section:
                    current_section['content'] += f"{line.strip()}\n"

        # 最后一部分内容需要添加到 sections 中
        if current_section:
            sections.append(current_section)

        return sections