from nltk import deprecated
from pptx.util import Inches


class ContentPPTCreator:
    def __init__(self, prs):
        """接收统一的 Presentation 对象"""
        self.prs = prs

    def add_content_slide(self, title, content):
        """添加内容页"""
        # 判断 content 是否包含 `###` 子标题
        if isinstance(content, str):  # 如果 content 是字符串
            lines = content.splitlines()  # 按行分割
        elif isinstance(content, list):  # 如果 content 是列表
            lines = content
        else:
            raise ValueError("Content must be a string or a list.")

        # 检查是否存在 `###` 开头的行
        has_triple_hash = any(line.strip().startswith("###") for line in lines)

        if has_triple_hash:
            # 如果存在 `###` 子标题，跳转到 add_multi_column_slide 函数
            return self.add_multi_column_slide(title, content)
        else:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
            # 设置标题
            title_shape = slide.shapes.title
            title_shape.text = title
            # 设置内容
            body_shape = slide.shapes.placeholders[1]
            if isinstance(content, list):  # 如果 content 是列表，将其合并为字符串
                content = "\n".join(content)
            body_shape.text = content

            return slide

    def add_two_column_slide(self, title, left_content, right_content):
        """添加双栏内容页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[3])

        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title

        # 左侧内容
        left_shape = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Inches(3.5), Inches(5)
        )
        left_tf = left_shape.text_frame
        left_tf.text = left_content

        # 右侧内容
        right_shape = slide.shapes.add_textbox(
            Inches(5), Inches(1.5),
            Inches(3.5), Inches(5)
        )
        right_tf = right_shape.text_frame
        right_tf.text = right_content

        return slide

    def add_multi_column_slide(self, title, content):
        """
        添加多栏内容页，将 `###` 子标题及其内容以分栏形式显示
        :param title: 幻灯片主标题
        :param content: 多栏内容（`###` 子标题及其正文）
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 添加空白布局幻灯片
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title

        # 获取幻灯片的总宽度和高度
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height

        # 定义边距
        margin_left = Inches(1)  # 左边距
        margin_top = title_shape.height + Inches(0.5)  # 标题下方的间距
        margin_right = Inches(1)  # 右边距

        # 可用宽度计算
        available_width = slide_width - margin_left - margin_right

        # 解析内容，将 `###` 子标题和后续内容分组
        grouped_content = []  # 存储分组的内容
        current_group = None  # 当前处理的分组

        # 按行解析内容
        lines = content.splitlines()
        for line in lines:
            line = line.strip()  # 去除空格
            if line.startswith("###"):  # 如果是子标题，创建新的组
                if current_group:  # 将现有组加入列表
                    grouped_content.append(current_group)
                current_group = {"subtitle": line[4:], "body": []}  # 初始化新组
            else:  # 如果是正文内容，加入当前组
                if current_group:
                    current_group["body"].append(line)

        # 将最后的组加入列表
        if current_group:
            grouped_content.append(current_group)

        # 动态计算列数（基于子标题数量）
        column_num = len(grouped_content)
        column_width = available_width / column_num  # 单列宽度

        # 按组创建文本框
        for i, group in enumerate(grouped_content):
            # 创建文本框
            left = margin_left + i * column_width  # 每列的左边距
            top = margin_top  # 每列的顶部边距

            textbox = slide.shapes.add_textbox(left, top, column_width, slide_height - margin_top)  # 添加文本框
            text_frame = textbox.text_frame
            text_frame.word_wrap = True  # 自动换行

            # 填充子标题
            subtitle_para = text_frame.add_paragraph()
            subtitle_para.text = group["subtitle"]  # 子标题文本
            subtitle_para.font.bold = True  # 子标题加粗
            subtitle_para.font.size = Inches(0.3)  # 字体大小设置

            # 填充正文内容
            for body_text in group["body"]:
                body_para = text_frame.add_paragraph()
                body_para.text = body_text
                body_para.font.size = Inches(0.25)  # 字体大小设置

        return slide

    def save(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)