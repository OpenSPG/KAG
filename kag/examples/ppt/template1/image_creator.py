import os
from io import BytesIO

import requests
from PIL import Image
from pptx.util import Inches


class ImagePPTCreator:
    def __init__(self, prs):
        """接收统一的 Presentation 对象"""
        self.prs = prs

    def add_image_slide_offline(self, title, image_path, figure_text=None, position=None):
        """添加一个包含离线图片的幻灯片，并附带标题和描述文字"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[9])

        if slide.shapes.title:
            # 如果存在标题占位符，设置标题文本
            title_shape = slide.shapes.title
            title_shape.text = title

        # 处理图片
        img_path = self.process_image(image_path)
        with Image.open(img_path) as img:
            img_width_px, img_height_px = img.size  # 宽度和高度，单位为像素
        # 设置图片位置
        if position is None:
            placeholder = slide.shapes[1]
            placeholder_height = placeholder.height  # 保持高度固定
            placeholder_left = placeholder.left
            placeholder_top = placeholder.top

            # 根据占位符高度，计算图片宽度，保持比例
            img_aspect_ratio = img_width_px / img_height_px  # 图片宽高比
            adjusted_width = int(placeholder_height * img_aspect_ratio)  # 根据高度调整宽度

            # 计算图片的水平居中位置
            adjusted_left = placeholder_left + (placeholder.width - adjusted_width) // 2

            # 添加图片到幻灯片
            slide.shapes.add_picture(
                img_path,
                adjusted_left,  # 居中位置
                placeholder_top,  # 占位符顶部位置
                adjusted_width,  # 动态调整的宽度
                placeholder_height,  # 固定高度
            )
            # 删除占位符
            slide.shapes._spTree.remove(placeholder._element)
        else:
            slide.shapes.add_picture(
                img_path,
                position.left,
                position.top,
                position.width,
                position.height,
            )

        # 添加图片的描述性文字（FigureText）
        if figure_text is not None:
            body_shape = slide.shapes.placeholders[2]
            body_shape.text = figure_text

        return slide

    def add_image_slide_online(self, title, image_url, figure_text=None, position=None):
        """添加一个包含在线图片的幻灯片，并附带标题和描述文字"""
        # 添加一个幻灯片
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[9])  # 使用空白布局
        # 添加标题
        if slide.shapes.title:
            # 如果存在标题占位符，设置标题文本
            title_shape = slide.shapes.title
            title_shape.text = title

        # 获取图片
        response = requests.get(image_url)
        if response.status_code == 200:
            image_data = BytesIO(response.content)
            # 设置图片位置
            if position is None:
                placeholder = slide.shapes[1]
                # 添加图片
                slide.shapes.add_picture(
                    image_data,
                    placeholder.left,
                    placeholder.top,
                    placeholder.width,
                    placeholder.height,
                )
                # 删除占位符
                slide.shapes._spTree.remove(placeholder._element)
            else:
                slide.shapes.add_picture(
                    image_data,
                    position.left,
                    position.top,
                    position.width,
                    position.height,
                )

        # 添加图片的描述性文字（FigureText）
        if figure_text is not None:
            body_shape = slide.shapes.placeholders[2]
            body_shape.text = figure_text

        return slide

    @staticmethod
    def process_image(image_path, max_size_mb=2):
        """处理图片尺寸和大小"""
        # 获取当前脚本文件的所在目录
        # absolute_image_path = os.path.join(os.getcwd(), image_path)
        img = Image.open(image_path)

        # 压缩图片直到小于指定大小
        while os.path.getsize(image_path) > max_size_mb * 1024 * 1024:
            width, height = img.size
            img = img.resize((int(width * 0.8), int(height * 0.8)))
            img.save(image_path)

        return image_path

    def save(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)