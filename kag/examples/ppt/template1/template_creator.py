from pptx import Presentation

from kag.examples.ppt.template1.ppt_creator_base import PPTCreatorBase


class TemplatePPTCreator(PPTCreatorBase):
    def apply_template(self, template_path, data):
        """应用模板"""
        self.prs = Presentation(template_path)

        # 遍历所有页面
        for slide_idx, slide in enumerate(self.prs.slides):
            # 遍历页面中的所有形状
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                # 替换文本中的占位符
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in data.items():
                            placeholder = f'{{{key}}}'
                            if placeholder in run.text:
                                run.text = run.text.replace(
                                    placeholder,
                                    str(value)
                                )

    def create_from_template(self, template_path, data_list):
        """批量使用模板创建PPT"""
        for item in data_list:
            try:
                # 应用模板
                self.apply_template(template_path, item)

                # 保存文件
                output_path = f"output/{item['title']}.pptx"
                self.save(output_path)
                print(f"创建成功：{output_path}")

            except Exception as e:
                print(f"处理 {item['title']} 时出错：{str(e)}")
                continue

    def save(self, filename=f"template.pptx"):
        """保存PPT文件"""
        self.prs.save(filename)
        print(f"创建成功：{filename}")


