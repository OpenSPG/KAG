from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


class PPTCreator:
    def __init__(self):
        self.prs = Presentation()

    def add_title_slide(self, title, subtitle=None):
        """添加标题页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        # slide = self.prs.slides[1]  # 选择现有幻灯片
        # slide = self.Slides(1)
        # slide.Copy()
        # title_shape = slide.shapes.title
        title_shape = slide.placeholders[0]
        subtitle_shape = slide.placeholders[1]

        # 设置标题
        title_shape.text = title
        # 字体大小
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        # 字体加粗
        title_shape.text_frame.paragraphs[0].font.bold = True
        # 字体对齐方式
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 设置副标题
        if subtitle:
            subtitle_shape.text = subtitle
            # subtitle_shape.text_frame.paragraphs[0].font.size = Pt(32)

        return slide

    def save(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)
