from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx import Presentation


class ChartPPTCreator:
    def __init__(self, prs):
        """接收统一的 Presentation 对象"""
        self.prs = prs

    def add_chart_slide(self, title, chart_data, chart_type=None):
        """添加图表页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[8])

        if slide.shapes.title:
            # 如果存在标题占位符，设置标题文本
            title_shape = slide.shapes.title
            title_shape.text = title
        else:
            # 如果不存在标题占位符，手动添加一个文本框作为标题
            print("当前幻灯片没有标题占位符，正在手动添加文本框...")
            left = Inches(1)
            top = Inches(0.5)
            width = Inches(8)
            height = Inches(1)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = title

        # 创建图表数据
        chart_table = CategoryChartData()

        # 添加类别
        categories = chart_data['categories']
        values = chart_data.get('series')
        series_name = chart_data.get('series_name')

        chart_table.categories = categories
        chart_table.add_series(series_name, values)

        # 设置图表类型
        if chart_type is None:
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

        # 添加图表，设置图表样式
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 1:  # 假设占位符索引为 1
                placeholder = shape
                chart_left, chart_top, chart_width, chart_height = (
                    placeholder.left,
                    placeholder.top,
                    placeholder.width,
                    placeholder.height,
                )
                chart = slide.shapes.add_chart(
                    chart_type,  # 图表类型：柱状图
                    chart_left, chart_top, chart_width, chart_height, chart_table
                ).chart

                # 设置图表标题
                chart.has_title = True
                chart.chart_title.text_frame.text = "季度销售数据"

        return slide

    @staticmethod
    def _style_chart(chart):
        """设置图表样式"""
        # 设置标题
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)

        # 设置图例
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        # 设置网格线
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(200, 200, 200)

    def add_multi_series_chart(self, title: str, content: str, data_dict):
        """添加多系列图表"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[8])

        chart_data = CategoryChartData()
        chart_data.categories = data_dict['categories']
        chart_name = data_dict['title']
        # 添加多个系列
        for series_name, values in data_dict['series'].items():
            chart_data.add_series(series_name, values)
        # 添加图表，设置图表样式
        if slide.shapes.title:
            # 如果存在标题占位符，设置标题文本
            title_shape = slide.shapes.title
            title_shape.text = title

        shape = slide.shapes[1]
        shape.text = content

        shape = slide.shapes[2]
        placeholder = shape
        chart_left, chart_top, chart_width, chart_height = (
            placeholder.left,
            placeholder.top,
            placeholder.width,
            placeholder.height,
        )
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,  # 图表类型：柱状图
            chart_left, chart_top, chart_width, chart_height, chart_data
        ).chart
        # 删除占位符
        slide.shapes._spTree.remove(placeholder._element)
        # 设置图表标题
        chart.has_title = True

        return slide

    def save(self, filename):
        """保存PPT文件"""
        self.prs.save(filename)